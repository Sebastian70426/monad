import os
import fitz


def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text_parts = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            text_parts.append(text)
    doc.close()
    return "\n\n".join(text_parts)


def extract_text_from_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
    return "\n".join(text_parts)


def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return f.read()


def extract_text(file_path, file_type):
    if file_type == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_type == '.pptx':
        return extract_text_from_pptx(file_path)
    elif file_type in ['.txt', '.md']:
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")