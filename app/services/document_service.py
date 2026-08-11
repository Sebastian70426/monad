import os
import fitz


# ===== 结构化提取（保留结构信息：页码/幻灯片号） =====

def extract_structured(file_path, file_type, filename):
    """提取文本并保留结构信息（页码/幻灯片号）。

    返回: [{"text": str, "page": int, "source": str}, ...]
    """
    if file_type == '.pdf':
        return _extract_pdf_structured(file_path, filename)
    elif file_type == '.pptx':
        return _extract_pptx_structured(file_path, filename)
    elif file_type in ['.txt', '.md']:
        return _extract_txt_structured(file_path, filename)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def _extract_pdf_structured(file_path, filename):
    """按页提取 PDF，每页一个 section"""
    doc = fitz.open(file_path)
    sections = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            sections.append({"text": text, "page": i + 1, "source": filename})
    doc.close()
    return sections


def _extract_pptx_structured(file_path, filename):
    """按幻灯片提取 PPTX，每张 slide 一个 section"""
    from pptx import Presentation
    prs = Presentation(file_path)
    sections = []
    for i, slide in enumerate(prs.slides):
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
        if text_parts:
            sections.append({
                "text": "\n".join(text_parts),
                "page": i + 1,
                "source": filename
            })
    return sections


def _extract_txt_structured(file_path, filename):
    """按空行分段提取 TXT/MD，每段一个 section"""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        content = f.read()
    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
    sections = []
    for i, para in enumerate(paragraphs):
        sections.append({"text": para, "page": i + 1, "source": filename})
    return sections


# ===== 图片文档（视觉模型提取） =====

IMAGE_EXTS = ['.jpg', '.jpeg', '.png', '.webp']


def extract_image_text(file_path, filename):
    """用多模态模型提取图片中的全部文字与图表内容。

    返回: [{"text": str, "page": 1, "source": filename}]
    """
    from services.llm_client import get_llm_client, image_to_data_url
    client = get_llm_client()
    data_url = image_to_data_url(file_path)
    prompt = ("你是 OCR 与图表识别助手。请完整提取图片中的全部文字内容（保留段落与标题层级，"
              "公式用纯文本如 F = m * a），并简要描述图中的图表、结构或示意。"
              "若图片几乎没有文字，请用中文描述图片内容。")
    text = client.chat(
        messages=[{"role": "user", "content": [
            {"type": "text", "text": prompt},
            {"type": "image_url", "image_url": {"url": data_url}},
        ]}],
        temperature=0.2,
        max_tokens=3000
    )
    return [{"text": text or "", "page": 1, "source": filename}]
